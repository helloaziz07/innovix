"""
Innovix — Project Plan Export Service

Exports generated project plans to Markdown and PDF formats.
Also provides a narration-ready text format for Sarvam AI TTS.

Uses WeasyPrint for PDF generation (already in requirements.txt).
"""

import json
import logging
from typing import Optional, List, Dict, Any

from google import genai
from pydantic import BaseModel, Field
from pptx import Presentation as PPTXPresentation
import io
from app.core.config import settings

logger = logging.getLogger(__name__)


def export_to_markdown(project: Dict[str, Any]) -> str:
    """
    Export a project plan to a formatted Markdown string.

    Args:
        project: Full project record from Supabase (including project_plan, tech_stack, etc.)

    Returns:
        Formatted Markdown string.
    """
    plan = project.get("project_plan", {})
    title = project.get("title", "Untitled Project")
    idea = project.get("idea_text", "")

    sections = []

    # Header
    sections.append(f"# {title}")
    sections.append("")
    sections.append(f"> {idea}")
    sections.append("")
    sections.append("---")
    sections.append("")

    # Problem Validation
    pv = plan.get("problem_validation", {})
    if pv:
        sections.append("## 🎯 Problem Validation")
        sections.append("")
        sections.append(pv.get("summary", ""))
        sections.append("")
        if pv.get("target_users"):
            sections.append("**Target Users:**")
            for user in pv["target_users"]:
                sections.append(f"- {user}")
            sections.append("")
        if pv.get("pain_points"):
            sections.append("**Pain Points:**")
            for pp in pv["pain_points"]:
                sections.append(f"- {pp}")
            sections.append("")
        if pv.get("market_size"):
            sections.append(f"**Market Size:** {pv['market_size']}")
            sections.append("")

    # Existing Solutions
    solutions = plan.get("existing_solutions", [])
    if solutions:
        sections.append("## 🔍 Existing Solutions")
        sections.append("")
        sections.append("| Solution | Description | Pros | Cons | Pricing |")
        sections.append("|---|---|---|---|---|")
        for sol in solutions:
            pros = ", ".join(sol.get("pros", [])[:3])
            cons = ", ".join(sol.get("cons", [])[:3])
            sections.append(
                f"| **{sol.get('name', '')}** | {sol.get('description', '')[:80]} | {pros} | {cons} | {sol.get('pricing', 'N/A')} |"
            )
        sections.append("")

    # Innovation Opportunities
    innovations = plan.get("innovation_opportunities", [])
    if innovations:
        sections.append("## 💡 Innovation Opportunities")
        sections.append("")
        for inn in innovations:
            impact_emoji = {"high": "🔴", "medium": "🟡", "low": "🟢"}.get(inn.get("impact", ""), "⚪")
            sections.append(f"### {impact_emoji} {inn.get('area', '')}")
            sections.append(f"{inn.get('description', '')}")
            sections.append(f"- **Impact:** {inn.get('impact', 'N/A')} | **Feasibility:** {inn.get('feasibility', 'N/A')}")
            sections.append("")

    # Tech Stack
    tech_stack = plan.get("tech_stack", project.get("tech_stack", []))
    if tech_stack and isinstance(tech_stack, list):
        sections.append("## 🛠️ Tech Stack")
        sections.append("")
        sections.append("| Layer | Technology | Justification |")
        sections.append("|---|---|---|")
        for tech in tech_stack:
            sections.append(
                f"| {tech.get('layer', '')} | **{tech.get('technology', '')}** | {tech.get('justification', '')[:100]} |"
            )
        sections.append("")

    # Architecture
    arch = plan.get("architecture", project.get("architecture", {}))
    if arch and isinstance(arch, dict):
        sections.append("## 🏗️ System Architecture")
        sections.append("")
        mermaid = arch.get("mermaid_diagram", "")
        if mermaid:
            sections.append("```mermaid")
            sections.append(mermaid)
            sections.append("```")
            sections.append("")
        components = arch.get("components", [])
        if components:
            sections.append("### Components")
            sections.append("")
            for comp in components:
                sections.append(f"- **{comp.get('name', '')}** ({comp.get('type', '')}): {comp.get('description', '')}")
            sections.append("")

    # Roadmap
    roadmap = plan.get("roadmap", [])
    if roadmap:
        sections.append("## 🗺️ Development Roadmap")
        sections.append("")
        for phase in roadmap:
            sections.append(f"### Phase {phase.get('phase', '?')}: {phase.get('name', '')}")
            sections.append(f"*Duration: {phase.get('duration_weeks', '?')} week(s)*")
            sections.append("")
            sections.append(phase.get("description", ""))
            sections.append("")
            milestones = phase.get("milestones", [])
            for ms in milestones:
                prio = {"critical": "🔴", "high": "🟠", "medium": "🟡", "low": "🟢"}.get(ms.get("priority", ""), "⚪")
                sections.append(f"- {prio} **{ms.get('name', '')}**")
                for d in ms.get("deliverables", []):
                    sections.append(f"  - {d}")
            sections.append("")

    # Timeline
    timeline = plan.get("timeline", [])
    if timeline and isinstance(timeline, list):
        sections.append("## 📅 Weekly Timeline")
        sections.append("")
        sections.append("| Week | Phase | Focus | Tasks |")
        sections.append("|---|---|---|---|")
        for week in timeline:
            tasks = ", ".join(week.get("tasks", [])[:4])
            sections.append(
                f"| Week {week.get('week', '?')} | Phase {week.get('phase', '?')} | {week.get('focus_area', '')} | {tasks} |"
            )
        sections.append("")

    # API & Datasets
    apis = plan.get("api_datasets", [])
    if apis:
        sections.append("## 🔌 Recommended APIs & Datasets")
        sections.append("")
        for api in apis:
            sections.append(f"- **{api.get('name', '')}** ({api.get('type', '')}) — {api.get('description', '')}")
            if api.get("url"):
                sections.append(f"  - URL: {api['url']}")
            sections.append(f"  - Pricing: {api.get('pricing', 'N/A')} | Relevance: {api.get('relevance', '')}")
        sections.append("")

    # GitHub Repos
    repos = plan.get("github_repos", [])
    if repos:
        sections.append("## 📦 Relevant GitHub Repositories")
        sections.append("")
        for repo in repos:
            sections.append(f"- **[{repo.get('name', '')}]({repo.get('url', '')})** ⭐ {repo.get('stars', '?')}")
            sections.append(f"  - {repo.get('description', '')}")
            sections.append(f"  - *Use case:* {repo.get('use_case', '')}")
        sections.append("")

    # Risks
    risks = plan.get("risks", [])
    if risks:
        sections.append("## ⚠️ Project Risks")
        sections.append("")
        for risk in risks:
            impact_emoji = {"high": "🔴", "medium": "🟡", "low": "🟢"}.get(risk.get("impact", ""), "⚪")
            sections.append(f"- {impact_emoji} **{risk.get('risk', '')}**")
            sections.append(f"  - Mitigation: {risk.get('mitigation', '')}")
        sections.append("")

    # Documentation / Proposal
    docs = plan.get("documentation", {})
    if docs:
        sections.append("## 📄 Project Proposal")
        sections.append("")
        if docs.get("tagline"):
            sections.append(f"*{docs['tagline']}*")
            sections.append("")
        if docs.get("proposal_outline"):
            sections.append(docs["proposal_outline"])
            sections.append("")

    # Footer
    sections.append("---")
    sections.append("*Generated by Innovix — AI-Powered Research & Innovation Copilot*")

    return "\n".join(sections)


def export_to_pdf(project: Dict[str, Any]) -> bytes:
    """
    Export a project plan to PDF bytes using WeasyPrint.

    Args:
        project: Full project record from Supabase.

    Returns:
        PDF content as bytes.
    """
    import markdown as md
    from xhtml2pdf import pisa
    import io

    # First generate markdown, then convert to HTML, then to PDF
    md_content = export_to_markdown(project)

    # Convert Markdown → HTML
    html_body = md.markdown(
        md_content,
        extensions=["tables", "fenced_code", "nl2br"],
    )

    # Wrap in styled HTML document
    html_doc = f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

        body {{
            font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
            color: #1a1a2e;
            line-height: 1.7;
            max-width: 800px;
            margin: 0 auto;
            padding: 40px 30px;
            font-size: 11pt;
        }}
        h1 {{
            color: #6c3ce9;
            font-size: 24pt;
            border-bottom: 3px solid #6c3ce9;
            padding-bottom: 10px;
            margin-bottom: 20px;
        }}
        h2 {{
            color: #2d2b55;
            font-size: 16pt;
            margin-top: 30px;
            border-bottom: 1px solid #e0e0e0;
            padding-bottom: 6px;
        }}
        h3 {{
            color: #4a4a6a;
            font-size: 13pt;
            margin-top: 20px;
        }}
        blockquote {{
            border-left: 4px solid #6c3ce9;
            padding: 8px 16px;
            background: #f8f6ff;
            margin: 10px 0;
            font-style: italic;
            color: #555;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 15px 0;
            font-size: 10pt;
        }}
        th {{
            background: #6c3ce9;
            color: white;
            padding: 8px 12px;
            text-align: left;
        }}
        td {{
            border: 1px solid #e0e0e0;
            padding: 8px 12px;
        }}
        tr:nth-child(even) {{
            background: #f9f9ff;
        }}
        code {{
            background: #f0ecff;
            padding: 2px 6px;
            border-radius: 4px;
            font-size: 10pt;
        }}
        pre {{
            background: #1a1a2e;
            color: #e0e0e0;
            padding: 16px;
            border-radius: 8px;
            overflow-x: auto;
            font-size: 10pt;
        }}
        hr {{
            border: none;
            border-top: 1px solid #e0e0e0;
            margin: 25px 0;
        }}
        ul, ol {{
            padding-left: 24px;
        }}
        li {{
            margin: 4px 0;
        }}
        a {{
            color: #6c3ce9;
            text-decoration: none;
        }}
    </style>
</head>
<body>
{html_body}
    </body>
    </html>
    """
    
    result = io.BytesIO()
    pdf = pisa.pisaDocument(io.BytesIO(html_doc.encode("UTF-8")), result)
    if pdf.err:
        raise Exception("Failed to generate PDF document.")
    return result.getvalue()


def get_narration_text(project: Dict[str, Any]) -> str:
    """
    Generate a concise narration-ready text for TTS (Sarvam AI).

    Produces a spoken summary of the project plan — shorter and more
    conversational than the full markdown export.

    Args:
        project: Full project record.

    Returns:
        Plain text suitable for TTS synthesis.
    """
    plan = project.get("project_plan", {})
    title = project.get("title", "Untitled Project")
    idea = project.get("idea_text", "")

    parts = []
    parts.append(f"Project: {title}.")
    parts.append(f"The idea is: {idea}")
    parts.append("")

    # Problem validation
    pv = plan.get("problem_validation", {})
    if pv.get("summary"):
        parts.append(f"Problem Validation: {pv['summary']}")
        parts.append("")

    # Innovation opportunities (top 2)
    innovations = plan.get("innovation_opportunities", [])[:2]
    if innovations:
        parts.append("Key Innovation Opportunities:")
        for inn in innovations:
            parts.append(f"  {inn.get('area', '')}: {inn.get('description', '')}")
        parts.append("")

    # Tech stack summary
    tech = plan.get("tech_stack", [])
    if tech and isinstance(tech, list):
        tech_names = [t.get("technology", "") for t in tech[:6]]
        parts.append(f"Recommended Tech Stack: {', '.join(tech_names)}.")
        parts.append("")

    # Roadmap summary
    roadmap = plan.get("roadmap", [])
    total_weeks = plan.get("total_weeks", "")
    mvp_week = plan.get("mvp_ready_by_week", "")
    if roadmap:
        parts.append(f"The project has {len(roadmap)} development phases spanning {total_weeks} weeks.")
        if mvp_week:
            parts.append(f"An MVP is expected by week {mvp_week}.")
        for phase in roadmap[:3]:
            parts.append(f"  Phase {phase.get('phase', '?')}: {phase.get('name', '')} — {phase.get('description', '')}")
        parts.append("")

    # Risks (top 2)
    risks = plan.get("risks", [])[:2]
    if risks:
        parts.append("Key Risks:")
        for r in risks:
            parts.append(f"  {r.get('risk', '')}. Mitigation: {r.get('mitigation', '')}")

    return "\n".join(parts)


class Slide(BaseModel):
    title: str = Field(description="The title of the slide")
    bullet_points: List[str] = Field(description="A list of bullet points for the slide content")

class Presentation(BaseModel):
    slides: List[Slide] = Field(description="List of slides in the presentation")

async def export_to_pptx(data: Dict[str, Any]) -> bytes:
    """
    Export project or workspace data to PPTX using Gemini to structure it.
    """
    client = genai.Client(api_key=settings.gemini_api_key)
    
    # Extract important parts instead of dumping entire data
    plan = data.get("project_plan", {}) if "project_plan" in data else data
    title = data.get("title", data.get("name", "Untitled Presentation"))
    
    summary = {
        "title": title,
        "idea": data.get("idea_text", ""),
        "problem_validation": plan.get("problem_validation", {}),
        "existing_solutions": plan.get("existing_solutions", []),
        "tech_stack": plan.get("tech_stack", []),
        "roadmap": plan.get("roadmap", []),
        "risks": plan.get("risks", [])
    }
    
    json_data = json.dumps(summary, indent=2, default=str)
    
    prompt = f"""
    You are an expert presentation creator. Summarize the following project/workspace data into a structured presentation format.
    Create a compelling narrative flow. Keep the bullet points concise and professional. Aim for 5-10 slides.
    Always include a Title slide and an Agenda slide at the beginning.
    
    Data:
    {json_data[:20000]}
    """
    
    try:
        response = client.models.generate_content(
            model="gemini-3.5-flash-lite",
            contents=prompt,
            config=genai.types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=Presentation,
            )
        )
        presentation_data = response.parsed
    except Exception as e:
        logger.error(f"Gemini generation failed: {e}")
        # Fallback to a basic structure if AI fails
        presentation_data = Presentation(slides=[
            Slide(title=title, bullet_points=["Project Overview", "Auto-generated slides"]),
            Slide(title="Error", bullet_points=["Failed to generate AI slides", str(e)])
        ])
    
    # Now generate PPTX
    prs = PPTXPresentation()
    for slide_data in presentation_data.slides:
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = slide_data.title
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, bp in enumerate(slide_data.bullet_points):
            if i == 0:
                tf.text = bp
            else:
                p = tf.add_paragraph()
                p.text = bp
                
    # Return as bytes
    result = io.BytesIO()
    prs.save(result)
    return result.getvalue()
